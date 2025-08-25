from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import ColorFormat
import requests
from io import BytesIO
import re

# Define color scheme
BG_COLOR = RGBColor(10, 10, 30)      # #0A0A1E - Dark navy
ACCENT_GREEN = RGBColor(0, 255, 136)  # #00FF88 - Neon green
ACCENT_BLUE = RGBColor(30, 144, 255)  # #1E90FF - Electric blue
ACCENT_PURPLE = RGBColor(107, 47, 255) # #6B2FFF - Deep purple
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(200, 200, 200)

def apply_background_theme(slide):
    """Apply consistent background theme to all slides"""
    # Main background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Accent sidebar
    sidebar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), Inches(0.2), Inches(7.5)
    )
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = ACCENT_BLUE
    sidebar.line.fill.background()
    
    # Top accent line
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.2), Inches(0), Inches(9.8), Inches(0.05)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = ACCENT_GREEN
    top_line.line.fill.background()

def add_page_number(slide, page_num):
    """Add page number footer to slide"""
    page_box = slide.shapes.add_textbox(Inches(8.5), Inches(7), Inches(1), Inches(0.3))
    page_frame = page_box.text_frame
    page_para = page_frame.paragraphs[0]
    page_para.text = str(page_num)
    page_para.font.color.rgb = GRAY
    page_para.font.size = Pt(12)
    page_para.alignment = PP_ALIGN.RIGHT

def add_title(prs, title, subtitle=""):
    """Add title slide with theme styling"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    apply_background_theme(slide)
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.margin_left = Inches(0.2)
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.color.rgb = ACCENT_GREEN
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle if provided
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.margin_left = Inches(0.2)
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.color.rgb = WHITE
        subtitle_para.font.size = Pt(28)
        subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Decorative accent
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(6), Inches(6), Inches(0.1))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_PURPLE
    accent_bar.line.fill.background()
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"Notes: Course introduction covering {title}"
    
    return slide

def add_toc(prs, items):
    """Add table of contents slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    apply_background_theme(slide)
    # Fixed: Calculate page number after adding slide
    add_page_number(slide, len(prs.slides))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Course Overview"
    title_para.font.color.rgb = ACCENT_GREEN
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # TOC items
    toc_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(5))
    toc_frame = toc_box.text_frame
    toc_frame.margin_left = Inches(0.2)
    
    for i, item in enumerate(items):
        if i == 0:
            p = toc_frame.paragraphs[0]
        else:
            p = toc_frame.add_paragraph()
        
        p.text = f"• {item}"
        p.font.color.rgb = WHITE
        p.font.size = Pt(20)
        p.space_after = Pt(8)
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Notes: Course structure overview with all modules listed"
    
    return slide

def add_divider(prs, big_title):
    """Add module divider slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    apply_background_theme(slide)
    # Fixed: Calculate page number after adding slide
    add_page_number(slide, len(prs.slides))
    
    # Large title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.5))
    title_frame = title_box.text_frame
    title_frame.margin_left = Inches(0.2)
    title_para = title_frame.paragraphs[0]
    title_para.text = big_title
    title_para.font.color.rgb = ACCENT_GREEN
    title_para.font.size = Pt(42)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Bold accent bar
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.2), Inches(7), Inches(0.2))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_BLUE
    accent_bar.line.fill.background()
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"Notes: Module introduction for {big_title}"
    
    return slide

def add_bullets(prs, title, bullets):
    """Add bullet slide with max 6 bullets, auto-split if needed"""
    bullet_chunks = [bullets[i:i+6] for i in range(0, len(bullets), 6)]
    
    slides_created = []
    for chunk_idx, bullet_chunk in enumerate(bullet_chunks):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        apply_background_theme(slide)
        # Fixed: Calculate page number after adding slide
        add_page_number(slide, len(prs.slides))
        
        # Title with part indicator if multiple slides
        slide_title = title
        if len(bullet_chunks) > 1:
            slide_title += f" (Part {chunk_idx + 1})"
            
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(1))
        title_frame = title_box.text_frame
        title_frame.margin_left = Inches(0.2)
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_title
        title_para.font.color.rgb = ACCENT_BLUE
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        
        # Bullet points
        bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8), Inches(5))
        bullet_frame = bullet_box.text_frame
        bullet_frame.margin_left = Inches(0.2)
        bullet_frame.margin_top = Inches(0.1)
        
        for i, bullet in enumerate(bullet_chunk):
            # Clean bullet text (remove markdown formatting)
            clean_bullet = re.sub(r'\*\*(.*?)\*\*', r'\1', bullet.strip('• '))
            
            if i == 0:
                p = bullet_frame.paragraphs[0]
            else:
                p = bullet_frame.add_paragraph()
            
            p.text = f"• {clean_bullet}"
            p.font.color.rgb = WHITE
            p.font.size = Pt(20)
            p.space_after = Pt(12)
            p.level = 0
        
        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = f"Notes: Key points for {title}"
        
        slides_created.append(slide)
    
    return slides_created

def add_table(prs, title, headers, rows):
    """Add table slide with neon styling"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    apply_background_theme(slide)
    # Fixed: Calculate page number after adding slide
    add_page_number(slide, len(prs.slides))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(1))
    title_frame = title_box.text_frame
    title_frame.margin_left = Inches(0.2)
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.color.rgb = ACCENT_BLUE
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    
    # Table
    table_rows = len(rows) + 1  # +1 for header
    table_cols = len(headers)
    
    table_shape = slide.shapes.add_table(
        table_rows, table_cols,
        Inches(0.8), Inches(2), Inches(8), Inches(4.5)
    )
    table = table_shape.table
    
    # Set equal column widths - Fixed: Convert to int
    col_width = int(Inches(8) / table_cols)
    for col_idx in range(table_cols):
        table.columns[col_idx].width = col_width
    
    # Header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_GREEN
        
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = BG_COLOR
            paragraph.font.bold = True
            paragraph.font.size = Pt(16)
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_data in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(25, 25, 50)  # Dark blue fill
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.color.rgb = WHITE
                paragraph.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"Notes: Data table for {title}"
    
    return slide

def add_image(prs, title, image_url, caption=""):
    """Add image slide with caption"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    apply_background_theme(slide)
    # Fixed: Calculate page number after adding slide
    add_page_number(slide, len(prs.slides))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(1))
    title_frame = title_box.text_frame
    title_frame.margin_left = Inches(0.2)
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.color.rgb = ACCENT_BLUE
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    
    # Image placeholder
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2), Inches(2), Inches(6), Inches(3.5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = RGBColor(40, 40, 60)
    img_placeholder.line.color.rgb = ACCENT_PURPLE
    img_placeholder.line.width = Pt(2)
    
    # Image placeholder text
    img_text = img_placeholder.text_frame
    img_para = img_text.paragraphs[0]
    img_para.text = f"IMAGE: {image_url.split('/')[-1]}"
    img_para.font.color.rgb = GRAY
    img_para.font.size = Pt(14)
    img_para.alignment = PP_ALIGN.CENTER
    
    # Caption
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(6), Inches(1))
        caption_frame = caption_box.text_frame
        caption_frame.margin_left = Inches(0.1)
        caption_para = caption_frame.paragraphs[0]
        caption_para.text = caption
        caption_para.font.color.rgb = ACCENT_PURPLE
        caption_para.font.size = Pt(16)
        caption_para.font.italic = True
        caption_para.alignment = PP_ALIGN.CENTER
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"Notes: Visual example for {title}"
    
    return slide

def add_code(prs, title, code_text):
    """Add code slide with monospace text"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    apply_background_theme(slide)
    # Fixed: Calculate page number after adding slide
    add_page_number(slide, len(prs.slides))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(1))
    title_frame = title_box.text_frame
    title_frame.margin_left = Inches(0.2)
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.color.rgb = ACCENT_BLUE
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    
    # Code box
    code_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8), Inches(5))
    code_frame = code_box.text_frame
    code_frame.margin_left = Inches(0.2)
    code_para = code_frame.paragraphs[0]
    code_para.text = code_text
    code_para.font.color.rgb = ACCENT_GREEN
    code_para.font.size = Pt(14)
    code_para.font.name = "Courier New"  # Monospace font
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"Notes: Code example for {title}"
    
    return slide

def add_demo_slide(prs, title, demo_info):
    """Add demo placeholder slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    apply_background_theme(slide)
    # Fixed: Calculate page number after adding slide
    add_page_number(slide, len(prs.slides))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(1))
    title_frame = title_box.text_frame
    title_frame.margin_left = Inches(0.2)
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.color.rgb = ACCENT_BLUE
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    
    # Demo area
    demo_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2), Inches(2), Inches(6), Inches(4)
    )
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = RGBColor(20, 20, 40)
    demo_box.line.color.rgb = ACCENT_GREEN
    demo_box.line.width = Pt(3)
    
    demo_text = demo_box.text_frame
    demo_para = demo_text.paragraphs[0]
    demo_para.text = "DEMO SPACE\n\nScreen Recording Insert Here"
    demo_para.font.color.rgb = WHITE
    demo_para.font.size = Pt(24)
    demo_para.alignment = PP_ALIGN.CENTER
    
    # Demo info
    info_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
    info_frame = info_box.text_frame
    info_para = info_frame.paragraphs[0]
    info_para.text = demo_info
    info_para.font.color.rgb = ACCENT_PURPLE
    info_para.font.size = Pt(16)
    info_para.alignment = PP_ALIGN.CENTER
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"Notes: Live demonstration for {title}"
    
    return slide

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    
    # Demo content structure
    content_structure = {
        'intro': {
            'title': 'Welcome to AI Entrepreneurship',
            'slides': [
                {
                    'title': 'Welcome to Your AI Business Journey',
                    'subtitle': 'From Zero to $10K+ Monthly in AI Services',
                    'type': 'content',
                    'content': [
                        '10 practical modules teaching profitable AI business skills',
                        'Real-world applications that businesses pay premium rates for',
                        'Free tools creating enterprise-level service capabilities',
                        'Systematic approach to building sustainable income streams'
                    ]
                },
                {
                    'title': 'The AI Opportunity Window',
                    'subtitle': 'The Perfect Storm for AI Entrepreneurs',
                    'type': 'content',
                    'content': [
                        'Every business needs AI integration but lacks internal expertise',
                        'Traditional consultants charge premium rates for basic implementations',
                        'Free AI tools now rival expensive enterprise software',
                        'Market timing - early adoption phase with massive opportunity'
                    ]
                }
            ]
        },
        'module_1': {
            'title': 'Module 1: AI Content Creation for Freelancers',
            'slides': [
                {
                    'title': 'The Content Creation Gold Rush',
                    'type': 'content',
                    'content': [
                        'Every business needs fresh, engaging content to attract customers',
                        'Traditional agencies charge premium rates for essential but difficult work',
                        'Content demand has exploded across all platforms',
                        'AI gives you a competitive advantage in this growing market'
                    ]
                },
                {
                    'title': 'Service Offerings & Pricing',
                    'type': 'table',
                    'table_data': [
                        ['Service Type', 'Price Range', 'Details'],
                        ['Blog Posts', '$25-$100', 'Per article, varies by length'],
                        ['Social Media Packages', '$200-$500/month', 'Multiple platforms'],
                        ['Email Newsletters', '$50-$150', 'Per newsletter'],
                        ['Product Descriptions', '$5-$15', 'Per item, volume adds up']
                    ]
                },
                {
                    'title': 'Live Demo - Coffee Shop Content Package',
                    'type': 'demo'
                }
            ]
        },
        'conclusion': {
            'title': 'Your AI Entrepreneurship Success Plan',
            'slides': [
                {
                    'title': 'Your Success Journey Starts Now',
                    'subtitle': 'From Learning to Earning Begins Today',
                    'type': 'content',
                    'content': [
                        'The market opportunity is real and waiting for someone to serve it',
                        'You have the knowledge to deliver genuine value to businesses',
                        'Action beats perfection in building successful enterprises',
                        'Your AI business journey begins with your next decision'
                    ]
                }
            ]
        }
    }
    
    # Create title slide
    add_title(prs, "AI Entrepreneurship Fundamentals", "Complete Course: From Zero to $10K+ Monthly")
    
    # Create table of contents
    toc_items = [
        "Introduction",
        "Module 1: AI Content Creation for Freelancers",
        "Module 2: Automated Lead Generation Systems", 
        "Module 3: AI-Powered Virtual Assistant Services",
        "Module 4: Poster & Flyer Design with Free AI Tools",
        "Module 5: Social Media Automation & Management",
        "Module 6: Data Analysis & Business Intelligence",
        "Module 7: E-commerce Automation & Optimization",
        "Module 8: Building Scalable AI Service Businesses",
        "Module 9: AI Video Content Creation & Monetization",
        "Module 10: Personal Brand & Thought Leadership with AI",
        "Conclusion"
    ]
    add_toc(prs, toc_items)
    
    # Process each section
    for section_key, section_data in content_structure.items():
        # Add section divider
        add_divider(prs, section_data['title'])
        
        # Process slides in section
        for slide_data in section_data['slides']:
            if slide_data['type'] == 'table' and slide_data.get('table_data'):
                headers = slide_data['table_data'][0]
                rows = slide_data['table_data'][1:]
                add_table(prs, slide_data.get('subtitle', slide_data['title']), headers, rows)
            
            elif slide_data['type'] == 'demo':
                add_demo_slide(prs, slide_data['title'], "Demo Duration: 3-4 minutes")
            
            elif slide_data.get('content'):
                add_bullets(prs, slide_data.get('subtitle', slide_data['title']), slide_data['content'])
    
    return prs

def main():
    """Main function to create and save presentation"""
    try:
        print("Creating AI Entrepreneurship Course presentation...")
        presentation = create_presentation()
        
        filename = "Course_AllModules.pptx"
        presentation.save(filename)
        
        print(f"✓ Successfully created {filename}")
        print(f"Total slides: {len(presentation.slides)}")
        print("Presentation includes:")
        print("- Title slide")
        print("- Table of contents")
        print("- Introduction section")
        print("- Module sections")
        print("- Conclusion section")
        print("- Professional dark theme with neon accents")
        print("- Page numbers and speaker notes")
        
    except Exception as e:
        print(f"Error creating presentation: {e}")

if __name__ == "__main__":
    main()